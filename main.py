import streamlit as st
import os
import json
import glob
import re
import hashlib
from collections import Counter

# --- IMPORTATIONS LANGCHAIN (Standards) ---
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_community.llms import Ollama
from langchain_classic.chains.combine_documents import create_stuff_documents_chain
from langchain_classic.chains import create_retrieval_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from pydantic import BaseModel, Field

# --- IMPORTATIONS RESTITUTION ---
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit_mermaid as st_mermaid

def clean_mermaid_code(text):
    """
    Nettoie le code Mermaid pour éviter les erreurs de syntaxe 10.2.4.
    """
    # 1. Supprimer la réflexion de DeepSeek (<think>...</think>)
    if "</think>" in text:
        text = text.split("</think>")[-1]

    # 2. Chercher le bloc de code Mermaid avec une Regex (plus précis)
    pattern = r"```mermaid\s*(.*?)\s*```"
    match = re.search(pattern, text, re.DOTALL)
    
    if match:
        code = match.group(1).strip()
    else:
        # Si pas de balises, on essaie de nettoyer le texte brut
        code = text.replace("```mermaid", "").replace("```", "").strip()

    # 3. Validation et Correction de secours
    lines = code.split('\n')
    clean_lines = []
    
    for line in lines:
        line = line.strip()
        # On ne garde que les lignes qui ressemblent à du code Mermaid ou des commentaires
        if line and (line.startswith("graph") or "-->" in line or line.startswith("subgraph") or line.startswith("end")):
            clean_lines.append(line)
            
    final_code = "\n".join(clean_lines)

    # Si l'IA a oublié de déclarer le type de graphique, on le force
    if not final_code.startswith("graph"):
        final_code = "graph TD\n" + final_code
        
    return final_code


# --- CONFIGURATION DE LA PAGE ---
st.set_page_config(page_title="Spaceflight Institute", page_icon="🚀", layout="wide")
st.title("🤖 Spaceflight Institute - Recherche Intelligente")

os.environ["NO_PROXY"] = "localhost,127.0.0.1"

# --- GESTION DES DOSSIERS ---
base_folder = "data"
cours_folder = os.path.join(base_folder, "cours")
users_folder = os.path.join(base_folder, "users")

for folder in [base_folder, cours_folder, users_folder]:
    if not os.path.exists(folder):
        os.makedirs(folder)

# --- 1. FONCTIONS AUTHENTIFICATION (Identique précédent) ---
def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def get_user_filepath(username):
    safe_name = re.sub(r'[^a-z0-9]', '', username.lower())
    return os.path.join(users_folder, f"{safe_name}.json")

# --- DANS LA SECTION 1. FONCTIONS AUTHENTIFICATION ---

def create_user(username, password, level, tone, role, goal):
    filepath = get_user_filepath(username)
    if os.path.exists(filepath): return False, "Utilisateur existant."
    
    data = {
        "auth": {"username_display": username, "password_hash": hash_password(password)},
        "profil": {
            "niveau": level,
            "role": role,           # NOUVEAU
            "objectif": goal,       # NOUVEAU
            "preferences_apprentissage": {"ton": tone}
        }
    }
    try:
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
        return True, "Succès"
    except: return False, "Erreur écriture"

def update_user_profile(username, new_level, new_tone, new_role, new_goal):
    """Met à jour le profil complet."""
    filepath = get_user_filepath(username)
    if not os.path.exists(filepath):
        return None
    
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
            
        # Modification des champs
        data["profil"]["niveau"] = new_level
        data["profil"]["preferences_apprentissage"]["ton"] = new_tone
        
        # Ajout/Modif des nouveaux champs
        data["profil"]["role"] = new_role
        data["profil"]["objectif"] = new_goal
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
            
        return data
    except Exception as e:
        print(f"Erreur update: {e}")
        return None

def verify_credentials(username, password):
    filepath = get_user_filepath(username)
    if not os.path.exists(filepath): return None, "Inconnu"
    try:
        with open(filepath, 'r') as f: data = json.load(f)
        if data["auth"]["password_hash"] == hash_password(password): return data, "Succès"
        return None, "Mot de passe faux"
    except: return None, "Erreur fichier"

# --- 2. FONCTIONS FICHIERS & CONTENU ---

def get_relevant_files(prompt, pdf_folder_path):
    """
    Tente de trouver des fichiers par nom. 
    Retourne TOUS les fichiers si pas de correspondance précise (pour laisser le RAG chercher dedans).
    """
    all_pdfs = glob.glob(os.path.join(pdf_folder_path, "*.pdf"))
    if not prompt or not all_pdfs: return all_pdfs, True 

    mots_vides = ["le", "la", "les", "de", "du", "des", "un", "une", "fichier", "cours", "pdf", "sur"]
    cleaned_prompt = re.sub(r'[^\w\s]', '', prompt.lower())
    keywords = [word for word in cleaned_prompt.split() if word not in mots_vides and len(word) > 2]
    
    selected_files = []
    for pdf_path in all_pdfs:
        filename = os.path.basename(pdf_path).lower()
        if any(kw in filename for kw in keywords):
            selected_files.append(pdf_path)
            
    # Si on trouve des fichiers par nom, on est contents
    if selected_files:
        return list(set(selected_files)), False
    
    # Sinon, on renvoie TOUT pour que le RAG cherche le contenu "puits de potentiel" partout
    return all_pdfs, True


# 1. On définit la structure de réponse attendue (JSON strict)
class FileRequest(BaseModel):
    veut_telecharger: bool = Field(description="Vrai si l'utilisateur demande explicitement un document, un article, un pdf ou une source complète.")
    nom_fichier_cible: str = Field(description="Le nom exact du fichier PDF parmi la liste fournie qui correspond le mieux à la demande. Vide si aucun fichier ne correspond.")
    raisonnement: str = Field(description="Pourquoi ce fichier a été choisi.")

def smart_file_router(user_prompt, folder_path):
    """
    Utilise l'IA pour décider quel fichier proposer, basé sur le sens de la phrase
    plutôt que sur des mots-clés.
    """
    # 1. Récupérer la liste des fichiers réels
    all_pdfs = [os.path.basename(f) for f in glob.glob(os.path.join(folder_path, "*.pdf"))]
    
    if not all_pdfs:
        return None

    # 2. Préparer le LLM (On utilise le même modèle)
    llm = Ollama(model="deepseek-r1:8b", temperature=0) # Température 0 pour être logique et strict
    
    # 3. Le Prompt du "Bibliothécaire"
    router_prompt = (
        "Tu es un bibliothécaire intelligent gérant une base de données de fichiers PDF. "
        "Ta mission est UNIQUEMENT d'analyser si l'utilisateur veut récupérer un fichier spécifique.\n"
        "Voici la liste EXACTE des fichiers disponibles dans ta bibliothèque :\n"
        f"{json.dumps(all_pdfs)}\n\n"
        "Analyses la demande de l'utilisateur. Si l'utilisateur demande un document, un article, un cours, ou une référence "
        "qui semble correspondre au SUJET d'un des fichiers, tu dois l'identifier.\n"
        "Exemple : Si l'utilisateur demande 'L'article de Palessandro' et que tu as 'Physique_Review_2024.pdf', et que tu sais ou devines que c'est lié, tu le sélectionnes.\n"
        "Attention : Réponds UNIQUEMENT au format JSON strict."
    )
    
    # 4. Création de la chaîne
    parser = JsonOutputParser(pydantic_object=FileRequest)
    prompt_template = ChatPromptTemplate.from_messages([
        ("system", router_prompt),
        ("human", "{input}\n\nFormat JSON attendu:\n{format_instructions}")
    ])
    
    chain = prompt_template | llm | parser
    
    try:
        # Lancement de l'analyse
        result = chain.invoke({
            "input": user_prompt,
            "format_instructions": parser.get_format_instructions()
        })
        
        # 5. Vérification du résultat
        if result["veut_telecharger"] and result["nom_fichier_cible"] in all_pdfs:
            # On retourne le chemin complet
            return os.path.join(folder_path, result["nom_fichier_cible"])
            
    except Exception as e:
        print(f"Erreur du routeur intelligent : {e}")
        # En cas d'erreur (si le LLM hallucine le JSON), on ne fait rien par sécurité
        return None
        
    return None




# --- 3. FONCTIONS DE GÉNÉRATION DE FORMATS ---
def generate_audio(text, lang='fr'):
    """Génère un fichier MP3 à partir du texte"""
    try:
        tts = gTTS(text=text, lang=lang, slow=False)
        filename = "temp_audio.mp3"
        tts.save(filename)
        return filename
    except Exception as e:
        return None

def generate_pptx_from_json(slides_data):
    """Crée un fichier PowerPoint à partir d'une liste de dictionnaires"""
    prs = Presentation()
    
    # Titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Synthèse Générée par IA"
    subtitle.text = "Spaceflight Institute"

    # Contenu
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_content in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_content.get("titre", "Sans titre")
        tf = body_shape.text_frame
        
        points = slide_content.get("points", [])
        if points:
            tf.text = points[0]
            for point in points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    output_file = "synthese_cours.pptx"
    prs.save(output_file)
    return output_file

# --- PROMPTS SPÉCIAUX (Formats uniquement) ---
PROMPT_MODES = {
    "Chat Standard": "Réponds de manière naturelle et pédagogique.",
    
    "🎙️ Résumé Audio (Podcast)": (
        "Adopte le persona d'un animateur de podcast dynamique. "
        "Ton but est de créer un script oral engageant. "
        "Ne mets pas de balises de mise en scène (musique, applaudissements), juste le texte parlé."
        "Ne lis SURTOUT PAS la ponctuation."
    ),
    
    "🧠 Carte Mentale": (
        "Tu es un expert en visualisation de données. "
        "Ta mission : Créer une carte mentale hiérarchique avec Mermaid.js. "
        "RÈGLES DE SYNTAXE STRICTES :\n"
        "1. Utilise UNIQUEMENT la syntaxe 'graph TD' (et non 'mindmap').\n"
        "2. IMPORTANT : Mets TOUS les textes des nœuds entre guillemets doubles. Ex: id[\"Mon Texte Ici\"]\n"
        "3. Ne mets JAMAIS de parenthèses () directement dans le texte sans guillemets.\n"
        "4. Réponds UNIQUEMENT avec le bloc de code ```mermaid ... ```.\n\n"
        "Exemple valide :\n"
        "```mermaid\n"
        "graph TD\n"
        "  A[\"Sujet Principal\"] --> B[\"Idée 1 (Détail)\"]\n"
        "  A --> C[\"Idée 2 : L'exemple\"]\n"
        "```"
    ),
    
    "📝 Fiches de Révision": (
        "Ton but est de créer un outil de mémorisation. "
        "Tu DOIS répondre UNIQUEMENT au format JSON strict (liste d'objets). "
        "Format attendu : [{{'question': '...', 'reponse': '...'}}, {{'question': '...', 'reponse': '...'}}]"
    ),
    
    "📊 Diapositives (PPTX)": (
        "Ton but est de synthétiser pour une présentation. "
        "Tu DOIS répondre UNIQUEMENT au format JSON strict. "
        "Format attendu : [{{'titre': '...', 'points': ['...', '...']}}, ...]"
    )
}

def contextualize_question(input_question, chat_history):
    """
    Reformule la question utilisateur en intégrant le contexte de l'historique.
    Gère proprement le nettoyage des balises <think> de DeepSeek.
    """
    if not chat_history:
        return input_question
    
    # On ne garde que les 3 derniers échanges pour rester focus
    # (Trop d'historique peut embrouiller l'IA sur des petits modèles)
    last_exchanges = chat_history[-3:] 
    
    # Construction de l'historique sous forme de texte
    history_str = ""
    for msg in last_exchanges:
        role = "Utilisateur" if msg["role"] == "user" else "Assistant"
        content = msg["content"]
        # On évite d'inclure les gros blocs de raisonnement ou JSON dans l'historique pour économiser des tokens
        if len(content) > 200: 
            content = content[:200] + "..."
        history_str += f"{role}: {content}\n"
    
    llm = Ollama(model="deepseek-r1:8b", temperature=0.1) # Température basse pour être précis
    
    # Prompt optimisé pour la reformulation
    reformulation_prompt = (
        "Tu es un expert en linguistique. Ta tâche est de réécrire la dernière question de l'utilisateur "
        "pour qu'elle soit autonome et compréhensible sans l'historique de la conversation.\n\n"
        "--- HISTORIQUE ---\n"
        f"{history_str}\n"
        "------------------\n"
        f"Dernière question de l'Utilisateur : '{input_question}'\n\n"
        "CONSIGNES :\n"
        "1. Remplace les pronoms (il, elle, ça, le...) par les noms précis présents dans l'historique.\n"
        "2. Si la question est déjà claire, ne change rien.\n"
        "3. Réponds UNIQUEMENT par la question reformulée. Pas de politesse, pas de balises.\n"
    )
    
    try:
        response = llm.invoke(reformulation_prompt)
        
        # NETTOYAGE CRITIQUE : On enlève tout ce qui est pensée (<think>...</think>)
        if "</think>" in response:
            response = response.split("</think>")[-1]
            
        # On enlève les guillemets éventuels et les espaces
        reformulated = response.replace('"', '').strip()
        
        return reformulated
    except Exception as e:
        print(f"Erreur reformulation : {e}")
        return input_question

# --- 3. INITIALISATION RAG (Avec Citations Précises) ---
def initialize_rag_chain_dynamic(selected_files, user_data, custom_prompt=None):
    
    # 1. Infos utilisateur
    # 1. Chargement des données 
    profil = user_data.get("profil", {})
    user_name = user_data.get("auth", {}).get("username_display", "Étudiant")
    
    user_level = profil.get("niveau", "Intermédiaire")
    user_role = profil.get("role", "Cadet Spatial")       
    user_goal = profil.get("objectif", "Apprentissage")  
    ai_tone = profil.get("preferences_apprentissage", {}).get("ton", "neutre")

    # 2. Chargement des PDF & Nettoyage Métadonnées
    all_pages = []
    for pdf_path in selected_files:
        try:
            loader = PyPDFLoader(pdf_path)
            docs = loader.load()
            
            # --- NETTOYAGE DES SOURCES ---
            # On remplace le chemin complet par juste le nom du fichier
            # On ajoute +1 aux pages car PyPDF compte à partir de 0
            file_name = os.path.basename(pdf_path)
            for doc in docs:
                doc.metadata["source"] = file_name
                if "page" in doc.metadata:
                    doc.metadata["page"] += 1 
            
            all_pages.extend(docs)
        except: pass

    if not all_pages: return None

    # 3. Vectorisation
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    chunks = text_splitter.split_documents(all_pages)
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    vectorstore = Chroma.from_documents(documents=chunks, embedding=embeddings)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})
    llm = Ollama(model="deepseek-r1:8b", temperature=0)
    llm = Ollama(model="deepseek-r1:8b")
    
    # 4. CRÉATION DU FORMAT DE DOCUMENT (La clé pour les citations !)
    # Au lieu de donner juste le texte, on injecte la source et la page avant chaque extrait
    doc_prompt = PromptTemplate.from_template(
        "Source: {source} (Page {page})\nContenu: {page_content}"
    )

    # 5. CONSTRUCTION DU PROMPT SYSTÈME
    format_instruction = custom_prompt if custom_prompt else "Réponds normalement."
    
    system_prompt = (
        f"Tu es un assistant expert pour le Spaceflight Institute. "
        f"Ton interlocuteur est : {user_name}, qui a le grade de {user_role}. "
        f"Son niveau de connaissances est : {user_level}. "
        f"Son objectif actuel est : {user_goal}. "
        f"Ton style de réponse doit être : {ai_tone}. "
        "Utilise le contexte fourni pour répondre avec précision."
        "\n\nContexte:\n{context}"
        
        "--- INSTRUCTION DE FORMAT ---\n"
        f"{format_instruction}\n"
        "------------------------------\n\n"
        
        "--- LOGIQUE DE RÉPONSE ET CITATIONS ---\n"
        "Tu disposes d'extraits de documents ci-dessous. Chaque extrait commence par sa Source et son Numéro de Page.\n"
        "1. Si tu utilises une information du contexte, tu DOIS citer ta source à la fin de la phrase ou du paragraphe.\n"
        "   Format obligatoire : [Fichier : Nom.pdf, Page : X]\n"
        "2. Si l'information n'est pas dans le contexte, utilise tes connaissances mais précise : [Source : Connaissances Générales].\n"
        "---------------------------------------\n\n"
        
        "--- CONTEXTE DOCUMENTS ---\n"
        "{context}"
    )
    
    prompt_template = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("human", "{input}")
    ])
    
    # On passe le 'document_prompt' ici pour que {context} contienne bien les sources
    chain = create_stuff_documents_chain(llm, prompt_template, document_prompt=doc_prompt)
    rag = create_retrieval_chain(retriever, chain)
    return rag

# --- GESTION SESSION ---
if "user_session" not in st.session_state: st.session_state["user_session"] = None

# --- SIDEBAR (LOGIN & PROFIL) ---
with st.sidebar:
    st.header("🔒 Authentification")
    
    # CAS 1 : UTILISATEUR CONNECTÉ
    if st.session_state["user_session"]:
        user_info = st.session_state["user_session"]
        username = user_info['auth']['username_display']
        
        # Récupération du rôle pour l'affichage (ex: Astronaute Nayel)
        display_role = user_info["profil"].get("role", "Astronaute")
        st.success(f"{display_role} : **{username}**")
        
        # --- BLOC UNIQUE : MODIFIER MON PROFIL ---
        with st.expander("👤 Dossier Personnel"):
            st.caption("Mettez à jour vos accréditations.")
            
            # --- 1. RECUPERATION DONNEES ACTUELLES (Avec Fallback) ---
            raw_level = user_info["profil"].get("niveau", "Intermédiaire")
            raw_tone = user_info["profil"].get("preferences_apprentissage", {}).get("ton", "Neutre")
            raw_role = user_info["profil"].get("role", "Cadet")
            raw_goal = user_info["profil"].get("objectif", "")

            # --- 2. LISTES D'OPTIONS ---
            list_niveaux = ["Débutant", "Intermédiaire", "Avancé", "Expert"]
            list_styles = ["Cool", "Strict", "Neutre", "Scientifique", "Militaire"]
            list_roles = ["Cadet", "Pilote", "Ingénieur", "Scientifique", "Commandant", "Touriste"]

            # --- 3. LOGIQUE DE SECURITE (Pour anciens profils) ---
            # Niveau
            if raw_level not in list_niveaux:
                mapping_niveau = {"A": "Débutant", "B": "Intermédiaire", "C": "Avancé"}
                default_level = mapping_niveau.get(raw_level, "Intermédiaire")
            else:
                default_level = raw_level
            
            # Style
            default_tone = raw_tone if raw_tone in list_styles else "Neutre"
            
            # Rôle
            default_role = raw_role if raw_role in list_roles else "Cadet"

            # --- 4. FORMULAIRE ---
            with st.form("update_profile_full"):
                st.markdown("**Identité**")
                new_role = st.selectbox("Spécialité", list_roles, index=list_roles.index(default_role))
                new_goal = st.text_input("Objectif de Mission", value=raw_goal, placeholder="Ex: Certification Moteur Raptor")
                
                st.markdown("**Préférences IA**")
                new_level = st.select_slider("Niveau d'expertise", options=list_niveaux, value=default_level)
                new_tone = st.selectbox("Ton de l'IA", list_styles, index=list_styles.index(default_tone))
                
                if st.form_submit_button("💾 Mettre à jour le dossier"):
                    updated_data = update_user_profile(username, new_level, new_tone, new_role, new_goal)
                    if updated_data:
                        st.session_state["user_session"] = updated_data
                        st.toast("Dossier mis à jour !", icon="✅")
                        st.rerun()
                    else:
                        st.error("Erreur mise à jour.")

        st.divider()
        
        # --- SÉLECTION DU MODE ---
        st.header("⚙️ Mode de Sortie")
        selected_mode = st.radio(
            "Format de la réponse :",
            list(PROMPT_MODES.keys()),
            key="mode_radio_main" 
        )
        
        st.divider()
        
        if st.button("Se déconnecter", key="logout_btn_main"):
            st.session_state["user_session"] = None
            st.rerun()
            
        st.divider()
        
        uploaded = st.file_uploader("Ajouter PDF", type="pdf", key="pdf_uploader_main")
        if uploaded:
            save_path = os.path.join(cours_folder, uploaded.name)
            with open(save_path, "wb") as f: 
                f.write(uploaded.getbuffer())
            st.success("Document archivé !")

    # CAS 2 : UTILISATEUR NON CONNECTÉ
    else:
        t1, t2 = st.tabs(["Connexion", "Inscription"])
        
        with t1:
            with st.form("login_form"):
                u = st.text_input("Matricule (User)")
                p = st.text_input("Code d'accès (Pass)", type="password")
                if st.form_submit_button("S'identifier"):
                    d, m = verify_credentials(u, p)
                    if d: 
                        st.session_state["user_session"] = d
                        st.rerun()
                    else: 
                        st.error(m)
        
        with t2:
            st.markdown("Rejoindre le **Spaceflight Institute**")
            with st.form("signup_form"):
                u = st.text_input("Nouveau Matricule")
                p = st.text_input("Créer Code d'accès", type="password")
                
                c1, c2 = st.columns(2)
                with c1:
                    role_signup = st.selectbox("Spécialité", ["Cadet", "Pilote", "Ingénieur", "Scientifique", "Commandant", "Touriste"])
                with c2:
                    level_signup = st.selectbox("Niveau", ["Débutant", "Intermédiaire", "Avancé", "Expert"])
                
                goal_signup = st.text_input("Objectif Principal", placeholder="Ex: Apprendre la physique orbitale")
                tone_signup = st.selectbox("Style de l'assistant", ["Cool", "Strict", "Neutre", "Scientifique", "Militaire"])
                
                if st.form_submit_button("Initialiser le profil"):
                    created, msg = create_user(u, p, level_signup, tone_signup, role_signup, goal_signup)
                    if created:
                        st.success("Profil créé ! Connectez-vous.")
                    else:
                        st.error(msg)
                        
# --- ZONE PRINCIPALE ---
if not st.session_state["user_session"]: st.stop()

if "messages" not in st.session_state: st.session_state.messages = []
for m in st.session_state.messages:
    with st.chat_message(m["role"]): st.markdown(m["content"])

# --- ZONE DE CHAT (Finale) ---

if prompt := st.chat_input("Posez votre question..."):
    
    # A. Message Utilisateur
    with st.chat_message("user"):
        st.markdown(prompt)
    st.session_state.messages.append({"role": "user", "content": prompt})

    # B. Réponse Assistant
    with st.chat_message("assistant"):
        
        # 1. GESTION DE LA MÉMOIRE (Reformulation)
        real_prompt = prompt
        
        # S'il y a un historique (plus que le message qu'on vient d'ajouter)
        if len(st.session_state.messages) > 1:
            with st.status("🧠 Analyse du contexte...", expanded=False) as status:
                # On prend l'historique AVANT la question actuelle
                history_context = st.session_state.messages[:-1]
                
                # Appel à la fonction de reformulation
                new_prompt = contextualize_question(prompt, history_context)
                
                # Si la reformulation est différente de la question de base, on l'utilise
                if new_prompt != prompt:
                    real_prompt = new_prompt
                    status.write(f"Question reformulée : {real_prompt}")
                    status.update(label="Contexte compris !", state="complete")
                else:
                    status.update(label="Pas de contexte nécessaire.", state="complete")

        # 2. ROUTEUR DE FICHIER (Sur la question reformulée ou originale ?)
        # Mieux vaut utiliser le prompt original pour le téléchargement "Donne moi le fichier..."
        found_file_path = smart_file_router(prompt, cours_folder)
        
        if found_file_path:
            filename = os.path.basename(found_file_path)
            st.success(f"📂 Document trouvé : **{filename}**")
            with open(found_file_path, "rb") as f:
                st.download_button(f"⬇️ Télécharger {filename}", f, file_name=filename)
            # On ajoute une note au prompt réel pour l'IA
            real_prompt += f" (Note : Tu as déjà proposé le fichier {filename}.)"

        # 3. RAG (Sur la question REFORMULÉE 'real_prompt')
        # C'est ici que ça change tout : on cherche "détails sur le rayonnement" au lieu de "détails"
        relevant_files, is_global = get_relevant_files(real_prompt, cours_folder)
        
        if relevant_files:
            with st.spinner(f"Rédaction (Mode {selected_mode})..."):
                try:
                    # Choix du mode
                    custom_instruction = PROMPT_MODES[selected_mode] if selected_mode != "Chat Standard" else None

                    # Initialisation
                    rag_chain = initialize_rag_chain_dynamic(
                        relevant_files, 
                        st.session_state["user_session"],
                        custom_prompt=custom_instruction
                    )
                    
                    if rag_chain:
                        # On envoie la question REFORMULÉE au RAG
                        response = rag_chain.invoke({"input": real_prompt})
                        raw_answer = response["answer"]

                        # Nettoyage
                        final_content = raw_answer
                        if "</think>" in raw_answer:
                            final_content = raw_answer.split("</think>")[-1].strip()
                        
                        final_content_clean = final_content.replace("```json", "").replace("```mermaid", "").replace("```", "").strip()

                        # Affichage selon mode
                        if selected_mode == "Chat Standard":
                            st.markdown(final_content)
                        
                        elif selected_mode == "🎙️ Résumé Audio (Podcast)":
                            st.markdown("### 🎙️ Podcast")
                            with st.expander("Script"): st.write(final_content)
                            audio_file = generate_audio(final_content)
                            if audio_file: st.audio(audio_file)
                        
                        elif selected_mode == "🧠 Carte Mentale":
                            st.markdown("### 🧠 Carte Mentale")
                            mermaid_code = clean_mermaid_code(raw_answer)
                            try:
                                st_mermaid.st_mermaid(mermaid_code, height="500px")
                            except:
                                st.code(mermaid_code)
                        
                        elif selected_mode == "📝 Fiches de Révision":
                            st.markdown("### 📝 Flashcards")
                            try:
                                flashcards = json.loads(final_content_clean)
                                cols = st.columns(2)
                                for i, card in enumerate(flashcards):
                                    with cols[i % 2]:
                                        with st.expander(f"❓ {card.get('question', 'Q')}", expanded=False):
                                            st.info(f"💡 {card.get('reponse', 'R')}")
                            except:
                                st.warning("Erreur JSON.")
                                st.write(final_content)

                        elif selected_mode == "📊 Diapositives (PPTX)":
                            st.markdown("### 📊 PowerPoint")
                            try:
                                slides_data = json.loads(final_content_clean)
                                pptx_file = generate_pptx_from_json(slides_data)
                                with open(pptx_file, "rb") as f:
                                    st.download_button("⬇️ Télécharger .pptx", f, "cours.pptx")
                                for slide in slides_data:
                                    st.markdown(f"**📺 {slide.get('titre', 'Slide')}**")
                                    for p in slide.get('points', []):
                                        st.markdown(f"- {p}")
                            except:
                                st.warning("Erreur PPTX.")
                                st.write(final_content)

                        # Sauvegarde historique (Réponse nette)
                        st.session_state.messages.append({"role": "assistant", "content": final_content})

                except Exception as e:
                    st.error(f"Erreur technique : {e}")