import re
import os
import streamlit as st  # Importation ajoutée pour le feedback visuel
from gtts import gTTS
from pptx import Presentation

def clean_mermaid_code(text):
    """
    Nettoie le code Mermaid pour éviter les erreurs de syntaxe.
    """
    # 1. Supprimer la réflexion de DeepSeek (<think>...</think>)
    if "</think>" in text:
        text = text.split("</think>")[-1]

    # 2. Chercher le bloc de code Mermaid avec une Regex
    pattern = r"```mermaid\s*(.*?)\s*```"
    match = re.search(pattern, text, re.DOTALL)
    
    if match:
        code = match.group(1).strip()
    else:
        code = text.replace("```mermaid", "").replace("```", "").strip()

    # 3. Validation et Correction de secours
    lines = code.split('\n')
    clean_lines = []
    
    for line in lines:
        line = line.strip()
        if line and (line.startswith("graph") or "-->" in line or line.startswith("subgraph") or line.startswith("end")):
            clean_lines.append(line)
            
    final_code = "\n".join(clean_lines)

    if not final_code.startswith("graph"):
        final_code = "graph TD\n" + final_code
        
    return final_code

import os
import re  # Import indispensable pour le nettoyage
from gtts import gTTS
import streamlit as st

def generate_audio(text, lang='fr'):
    """Génère un fichier audio en nettoyant les caractères Markdown."""
    try:
        # 1. Nettoyage du texte : on retire les astérisques et autres symboles Markdown
        # On remplace les * par rien du tout
        clean_text = re.sub(r'\*', '', text)
        # Optionnel : on peut aussi retirer les underscores _ ou les # des titres
        clean_text = re.sub(r'#+', '', clean_text)
        
        # 2. Configuration du dossier
        output_folder = "generated"
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
            
        filename = os.path.join(output_folder, "generated.mp3")
        
        # 3. Génération avec le texte NETTOYÉ
        tts = gTTS(text=clean_text, lang=lang, slow=False)
        tts.save(filename)
        
        st.toast("Fichier audio généré avec succès !", icon="🔊")
        return filename
        
    except Exception as e:
        st.error(f"Erreur lors de la génération audio : {e}")
        return None

def generate_pptx_from_json(slides_data):
    """Génère un PowerPoint et affiche un message de succès."""
    try:
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Synthèse IA"
        
        bullet_layout = prs.slide_layouts[1]
        for slide_content in slides_data:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = slide_content.get("titre", "Slide")
            tf = slide.shapes.placeholders[1].text_frame
            for point in slide_content.get("points", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
        output_file = "synthese.pptx"
        prs.save(output_file)
        
        # Feedback visuel ajouté
        st.success(f"✅ Présentation '{output_file}' prête !")
        
        return output_file
    except Exception as e:
        st.error(f"Erreur lors de la création du PPTX : {e}")
        return None